# app.py — SE Oil & Gas Autoforecasting

import os
import pathlib
import tempfile
from io import BytesIO

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import streamlit as st

from core import (
    load_header, load_production, fill_lateral_by_geo,
    preprocess, PreprocessConfig, forecast_all, ForecastConfig,
    plot_one_well, forecast_one_well, _train_rf,
    compute_eur_stats, probit_plot, eur_summary_table,
    modified_arps, D_LIM_DEFAULT
)

st.set_page_config(page_title="SE Tool", layout="wide")
LOGO_PATH = pathlib.Path(__file__).parent / "static" / "logo.png"

# ---------- Header ----------
cols = st.columns([0.12, 0.88])
with cols[0]:
    if LOGO_PATH.exists():
        st.image(str(LOGO_PATH), use_column_width=True)
with cols[1]:
    st.success("All Systems Working")
    st.title("SE Oil & Gas Autoforecasting")

# ================= Sidebar =================
with st.sidebar:
    st.header("Global Parameters")
    vertical_wells = st.checkbox(
        "Vertical Well Dataset",
        value=False,
        help="Disables lateral length normalization. No lateral length required — all wells are retained."
    )
    if vertical_wells:
        use_norm = False
        norm_len = 10_000   # unused placeholder
        st.info("Normalization disabled for vertical wells.")
    else:
        norm_len = st.number_input("Normalization Length (ft)", 1000, 30000, 10000, step=500)
        use_norm = st.checkbox("Apply Normalization", value=True)
    st.markdown("---")
    b_low  = st.number_input("b-factor Low", value=0.8, step=0.1, format="%.3f")
    b_high = st.number_input("b-factor High", value=1.2, step=0.1, format="%.3f")
    st.markdown("---")
    lat_col = st.text_input("Latitude column (optional QC)", value="Latitude")
    lon_col = st.text_input("Longitude column (optional QC)", value="Longitude")
    bin_decimals = st.number_input("Geo Bin Decimals", 0, 4, 2)
    st.markdown("---")
    if st.button("Reset All"):
        for k in list(st.session_state.keys()):
            del st.session_state[k]
        st.rerun()

def _fmt2(value) -> str:
    try:
        f = float(value)
        if np.isfinite(f): return f"{f:.2f}"
        return ""
    except: return "" if value is None else str(value)

def format_df_2dec(df: pd.DataFrame) -> pd.DataFrame:
    if df is None or df.empty: return df
    out = df.copy()
    for c in out.columns:
        out[c] = out[c].map(_fmt2)
    return out

def _phase_color(fluid: str) -> str:
    return {'oil': 'green', 'gas': 'red', 'water': 'blue'}[fluid.lower()]

def _phase_color_rgb(fluid: str):
    """Return ReportLab color object for PDF generation"""
    from reportlab.lib import colors
    color_map = {
        'oil': colors.HexColor('#2E7D32'),    # Green
        'gas': colors.HexColor('#D32F2F'),    # Red
        'water': colors.HexColor('#1976D2')   # Blue
    }
    return color_map.get(fluid.lower(), colors.grey)

def _save_fig(fig, dpi=220):
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    fig.savefig(tmp.name, dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    return tmp.name


def _make_outlier_fig(oneline_df, eur_col, fluid_name):
    """EUR vs Lateral Length scatter with IQR-fence outlier detection.
    Returns (fig, upper_fence, n_active, n_removed, median_eur) or None."""
    df = oneline_df.copy()
    df['_eur'] = pd.to_numeric(df[eur_col], errors='coerce')
    if 'LateralLength' in df.columns:
        df['_lat'] = pd.to_numeric(df['LateralLength'], errors='coerce')
    else:
        return None
    df = df.dropna(subset=['_eur', '_lat'])
    if len(df) < 3:
        return None

    q1, q3 = df['_eur'].quantile(0.25), df['_eur'].quantile(0.75)
    iqr = q3 - q1
    upper_fence = q3 + 1.5 * iqr
    lower_fence = q1 - 1.5 * iqr

    high_mask = df['_eur'] > upper_fence
    low_mask  = df['_eur'] < lower_fence
    ok_mask   = ~high_mask & ~low_mask
    unit = "MMcf" if fluid_name.lower() == "gas" else "Mbbl"

    fig, ax = plt.subplots(figsize=(7, 4))
    if ok_mask.any():
        ax.scatter(df.loc[ok_mask, '_lat'], df.loc[ok_mask, '_eur'],
                   color='#2C2C2C', s=40, zorder=3, label='Within range')
    if high_mask.any():
        ax.scatter(df.loc[high_mask, '_lat'], df.loc[high_mask, '_eur'],
                   color='#781F28', s=40, zorder=3, label='High EUR outlier (>1.5\u00d7 IQR)')
    if low_mask.any():
        ax.scatter(df.loc[low_mask, '_lat'], df.loc[low_mask, '_eur'],
                   color='#BBBBBB', s=40, zorder=3, label='Low EUR outlier (<1.5\u00d7 IQR)')
    ax.axhline(upper_fence, color='#781F28', linestyle='--', linewidth=0.8, alpha=0.6)
    ax.text(0.98, 0.97, f"Upper fence  {upper_fence:,.0f}",
            transform=ax.transAxes, ha='right', va='top', color='#781F28', fontsize=8)
    for spine in ax.spines.values():
        spine.set_edgecolor('#E8AABB')
        spine.set_linewidth(1.5)
    ax.set_facecolor('#FAFAFA')
    fig.patch.set_facecolor('white')
    ax.set_xlabel('Lateral Length (ft)', fontsize=9)
    ax.set_ylabel(f'EUR ({unit})', fontsize=9)
    ax.legend(fontsize=8, loc='upper right', bbox_to_anchor=(0.99, 0.88))
    ax.tick_params(labelsize=8)
    plt.tight_layout()

    n_active  = int(ok_mask.sum())
    n_removed = int((high_mask | low_mask).sum())
    n_flagged = int(high_mask.sum())
    median_eur = float(df['_eur'].median())
    return fig, upper_fence, n_active, n_removed, n_flagged, median_eur, unit

# ================= Upload =================
st.header("Upload & Prepare Data")
c1, c2 = st.columns(2)
with c1:
    header_file = st.file_uploader("Header CSV", type=["csv"], key="header_csv")
with c2:
    prod_file   = st.file_uploader("Production CSV", type=["csv"], key="prod_csv")

if st.button("Load / QC / Merge"):
    if not header_file or not prod_file:
        st.error("Please upload both Header and Production CSV files.")
    else:
        try:
            header_bytes = header_file.getvalue()
            header_df = load_header(BytesIO(header_bytes))
            raw_hdr   = pd.read_csv(BytesIO(header_bytes))
            for col in [lat_col, lon_col]:
                if col in raw_hdr.columns and col not in header_df.columns:
                    header_df[col] = raw_hdr[col]
            header_qc = fill_lateral_by_geo(
                header_df, lat_col=lat_col, lon_col=lon_col,
                lateral_col='LateralLength', decimals=int(bin_decimals)
            )
            prod_df = load_production(prod_file)
            st.session_state.header_qc = header_qc
            st.session_state.prod_df   = prod_df
            pp_cfg = PreprocessConfig(
                normalization_length=int(norm_len),
                use_normalization=bool(use_norm)
            )
            merged = preprocess(header_qc, prod_df, pp_cfg)
            if merged.empty:
                st.warning("No rows after preprocessing. Check your inputs.")
            else:
                st.session_state.merged = merged
                st.success(f"Merged {len(merged['WellID'].unique())} wells and {len(merged)} monthly rows.")
        except Exception as e:
            st.exception(e)

if "merged" in st.session_state:
    with st.expander("Preview: Header (QC'd)"):
        st.dataframe(format_df_2dec(st.session_state.header_qc.head(20)), use_container_width=True)
    with st.expander("Preview: Production"):
        st.dataframe(format_df_2dec(st.session_state.prod_df.head(20)), use_container_width=True)
    with st.expander("Preview: Merged"):
        st.dataframe(format_df_2dec(st.session_state.merged.head(20)), use_container_width=True)

st.markdown("---")

# ---------- Analytics ----------
def bfactor_analytics_figures(oneline: pd.DataFrame, fluid: str, eur_col: str):
    color = _phase_color(fluid)
    df = oneline.copy()
    df = df[pd.to_numeric(df.get('b'), errors='coerce').notna()]
    if df.empty:
        return None, None, None, pd.DataFrame([{"message":"no b values"}])

    b = pd.to_numeric(df['b'], errors='coerce').dropna().astype(float)
    P10, P50, P90 = np.percentile(b, [10, 50, 90])
    stats_rows = [
        ["count", float(b.size)], ["mean", float(b.mean())], ["median", P50],
        ["P10", P10], ["P90", P90]
    ]
    stats_df = pd.DataFrame(stats_rows, columns=[f"{fluid} b-factor statistics", "value"])

    fig_h, axh = plt.subplots(figsize=(6.6, 3.4))
    axh.hist(b.values, bins=30, color=color, alpha=0.85, edgecolor='none')
    axh.axvline(P50, color='black', linestyle='--', label=f"Median={P50:.2f}")
    axh.legend()
    hist_png = _save_fig(fig_h)

    fig_bx, axbx = plt.subplots(figsize=(6.6, 1.1))
    axbx.boxplot(b.values, vert=False, medianprops=dict(color='black'))
    axbx.set_yticks([])
    box_png = _save_fig(fig_bx)

    scatter_png = None
    if eur_col in df.columns:
        eur = pd.to_numeric(df[eur_col], errors='coerce')
        mask = eur > 0
        if mask.any():
            x = np.log10(eur[mask].astype(float).values)
            y = b[mask].astype(float).values
            fig_s, axs = plt.subplots(figsize=(6.6, 3.2))
            axs.scatter(10**x, y, s=16, color=color, alpha=0.85)
            axs.set_xscale('log')
            scatter_png = _save_fig(fig_s)
    return hist_png, box_png, scatter_png, stats_df

def _fit_arps_to_trace(t: np.ndarray, y: np.ndarray, d_lim: float = D_LIM_DEFAULT):
    """Fit a Modified Arps curve to an empirical trace; returns (qi, b, di)."""
    from scipy.optimize import minimize
    y = np.clip(y, 1e-6, None)
    qi0, b0, di0 = float(y[0]), 1.0, 0.10
    def loss(p):
        qi, b, di = p
        pred = modified_arps(qi, b, di, d_lim, t)
        return np.mean((np.log(pred + 1e-6) - np.log(y + 1e-6)) ** 2)
    bounds = [(y.max() * 0.1, y.max() * 3.0), (0.01, 2.0), (d_lim, 1.5)]
    res = minimize(loss, [qi0, b0, di0], method='L-BFGS-B', bounds=bounds)
    return tuple(res.x) if res.success else (qi0, b0, di0)


def build_type_curves_and_lines(monthly_df: pd.DataFrame, com: str, min_wells: int = 1,
                                oneline: pd.DataFrame = None):
    vol_col = f"Monthly_{com}_volume"
    if monthly_df.empty or vol_col not in monthly_df.columns:
        return pd.DataFrame(columns=['t','P10','P50','P90']), []

    # Historical lines for background well traces
    hist = monthly_df[monthly_df['Segment'] == 'Historical'][['WellID','Date',vol_col]].dropna().copy()
    hist = hist.sort_values(['WellID','Date'])
    hist['t'] = hist.groupby('WellID').cumcount() + 1
    lines = []
    for _, g in hist.groupby('WellID'):
        t = g['t'].to_numpy()
        y = np.clip(g[vol_col].to_numpy(), 1e-6, None)
        lines.append((t, y))

    # Build smooth type curves from per-well Arps parameters sorted by EUR.
    # Grouping by EUR percentile and taking median qi/b/di per group produces
    # pure Modified Arps curves with no drops or kinks.
    eur_col = {'oil': 'EUR (Mbbl)', 'gas': 'EUR (MMcf)', 'water': 'EUR (Mbbl water)'}[com.lower()]
    if (oneline is not None and not oneline.empty
            and eur_col in oneline.columns
            and all(c in oneline.columns for c in ['qi (per day)', 'b', 'di (per month)'])):
        df = oneline.copy()
        df['_eur'] = pd.to_numeric(df[eur_col], errors='coerce')
        df['_qi']  = pd.to_numeric(df['qi (per day)'], errors='coerce') * 30.4  # monthly
        df['_b']   = pd.to_numeric(df['b'], errors='coerce')
        df['_di']  = pd.to_numeric(df['di (per month)'], errors='coerce')
        df = df.dropna(subset=['_eur', '_qi', '_b', '_di'])
        df = df.sort_values('_eur', ascending=False).reset_index(drop=True)
        n = len(df)
        if n >= min_wells:
            t_out = np.arange(1, 601, dtype=float)
            smooth = {'t': t_out}
            p50_lo = max(0, int(n * 0.45))
            p50_hi = max(p50_lo + 1, int(n * 0.55))
            groups = {
                'P10': df.iloc[:max(1, int(n * 0.1))],
                'P50': df.iloc[p50_lo:p50_hi],
                'P90': df.iloc[max(0, int(n * 0.9)):],
            }
            # All curves share P90's b and di (the representative decline shape).
            # qi per group is a placeholder — _render_tw scales each curve to its
            # EUR target, so only the shape (b, di) matters here.
            b_shape  = float(groups['P90']['_b'].median())
            di_shape = float(groups['P90']['_di'].median())
            for pct, grp in groups.items():
                qi = float(grp['_qi'].median())
                smooth[pct] = modified_arps(qi, b_shape, di_shape, D_LIM_DEFAULT, t_out)
            return pd.DataFrame(smooth), lines

    # Fallback: empirical percentiles with Arps fitting (used when oneline unavailable)
    all_df = monthly_df[['WellID','Date',vol_col]].dropna().sort_values(['WellID','Date']).copy()
    all_df['t'] = all_df.groupby('WellID').cumcount() + 1
    counts = all_df.groupby('t')[vol_col].count()
    valid_t = counts[counts >= min_wells].index
    if valid_t.empty:
        return pd.DataFrame(columns=['t','P10','P50','P90']), lines

    emp = all_df[all_df['t'].isin(valid_t)].groupby('t')[vol_col].quantile([0.90, 0.50, 0.10]).unstack(level=1)
    emp.columns = ['P10', 'P50', 'P90']
    emp = emp.reset_index()
    t_emp = emp['t'].to_numpy(dtype=float)

    t_out = np.arange(1, int(t_emp.max()) + 1, dtype=float)
    smooth = {'t': t_out}
    for pct in ['P10', 'P50', 'P90']:
        qi, b, di = _fit_arps_to_trace(t_emp, emp[pct].to_numpy())
        smooth[pct] = modified_arps(qi, b, di, D_LIM_DEFAULT, t_out)

    return pd.DataFrame(smooth), lines

def plot_type_curves(curves, lines, fluid):
    color = _phase_color(fluid)
    fig, ax = plt.subplots(figsize=(9,5))
    if lines:
        for t, y in lines:
            ax.plot(t, y, color='gray', alpha=0.3, linewidth=0.5)
    if not curves.empty:
        ax.plot(curves['t'], curves['P50'], color=color, linewidth=2, label='P50')
        ax.plot(curves['t'], curves['P10'], color=color, linestyle='--', label='P10')
        ax.plot(curves['t'], curves['P90'], color=color, linestyle='--', label='P90')
    ax.set_yscale('log'); ax.set_ylim(bottom=1)
    ax.legend()
    return fig

# ================= Fluid Block =================
def fluid_block(fluid_name: str, eur_col: str, norm_col_for_models: str):
    st.header(f"{fluid_name} — Run & Analyze")
    disabled = "merged" not in st.session_state
    if st.button(f"Run {fluid_name} Forecast", disabled=disabled):
        merged = st.session_state.merged
        cfg = ForecastConfig(
            commodity=fluid_name.lower(),
            b_low=float(b_low), b_high=float(b_high), max_months=600
        )
        oneline, monthly, well_fcs = forecast_all(merged, cfg)
        st.session_state[f"{fluid_name}_oneline"] = oneline
        st.session_state[f"{fluid_name}_monthly"] = monthly
        st.session_state[f"{fluid_name}_well_fcs"] = well_fcs   # cache per-well forecasts
        # Reset analog selection so new forecast starts with all wells selected
        st.session_state.pop(f"{fluid_name}_tc_selection", None)
        st.success(f"{fluid_name} forecast completed.")

    on_key = f"{fluid_name}_oneline"
    mo_key = f"{fluid_name}_monthly"

    if on_key in st.session_state:
        tab1, tab2, tab3, tab4 = st.tabs([
            "Oneline", "Monthly", "B-Factor", "Probit"
        ])
        with tab1: st.dataframe(format_df_2dec(st.session_state[on_key]), use_container_width=True)
        with tab2: st.dataframe(format_df_2dec(st.session_state[mo_key]), use_container_width=True)
        with tab3:
            oneline = st.session_state[on_key]
            hist_png, box_png, scatter_png, bstats = bfactor_analytics_figures(oneline, fluid_name, eur_col)
            c1, c2 = st.columns([2,1])
            with c1:
                if hist_png: st.image(hist_png)
                if scatter_png: st.image(scatter_png)
            with c2:
                if box_png: st.image(box_png)
                st.dataframe(format_df_2dec(bstats))
        with tab4:
            oneline = st.session_state[on_key]
            if eur_col in oneline.columns:
                eurs = pd.to_numeric(oneline[eur_col], errors="coerce").astype(float).tolist()
                unit = "Mbbl" if "Mbbl" in eur_col else "MMcf"
                stats = compute_eur_stats(eurs)
                st.dataframe(format_df_2dec(eur_summary_table(fluid_name, stats, unit, int(norm_len))))
                fig = probit_plot(eurs, unit, f"{fluid_name} Probit", _phase_color(fluid_name), norm_len=int(norm_len))
                st.pyplot(fig)

        st.subheader(f"{fluid_name} — Per-well Plot")
        merged = st.session_state.merged
        # Use WellID for selection
        opts = merged[['WellID']].drop_duplicates().sort_values('WellID')
        pick_id = st.selectbox(f"Pick Well ({fluid_name})", opts['WellID'], key=f"{fluid_name}_pick")
        
        wd = merged[merged['WellID'] == pick_id]
        fcs_key = f"{fluid_name}_well_fcs"
        if fcs_key in st.session_state and pick_id in st.session_state[fcs_key]:
            # Reuse cached forecast — no RF re-training or re-optimisation
            fc = st.session_state[fcs_key][pick_id]
        else:
            models = _train_rf(merged, norm_col_for_models)
            fc = forecast_one_well(wd, fluid_name.lower(), float(b_low), float(b_high), 600, models)
        fig = plot_one_well(wd, fc, fluid_name.lower())
        st.pyplot(fig)

fluid_block("Oil", "EUR (Mbbl)", "NormOil")
st.markdown("---")
fluid_block("Gas", "EUR (MMcf)", "NormGas")
st.markdown("---")
fluid_block("Water", "EUR (Mbbl water)", "NormWater")
st.markdown("---")

# ================= Outlier Detection =================
st.header("Outlier Detection — EUR vs Lateral Length")

_out_fluid_cfg = [
    ("Oil",   "Oil_oneline",   "EUR (Mbbl)"),
    ("Gas",   "Gas_oneline",   "EUR (MMcf)"),
    ("Water", "Water_oneline", "EUR (Mbbl water)"),
]

def _outlier_section(fluid: str, on_key: str, eur_col: str):
    if on_key not in st.session_state:
        st.info(f"Run {fluid} forecast first.")
        return

    oneline = st.session_state[on_key]
    result  = _make_outlier_fig(oneline, eur_col, fluid)
    if result is None:
        st.info(f"Not enough data for {fluid} outlier detection (need ≥ 3 wells with lateral length).")
        return

    fig, upper_fence, n_active, n_removed, n_flagged, median_eur, unit = result
    st.pyplot(fig)
    plt.close(fig)

    m1, m2, m3, m4, m5 = st.columns(5)
    m1.metric("Active wells",  n_active)
    m2.metric("Removed wells", n_removed)
    m3.metric("Median EUR",    f"{median_eur:,.0f} {unit}")
    m4.metric("Upper fence",   f"{upper_fence:,.0f} {unit}")
    m5.metric("Flagged",       f"{n_flagged} wells")

    # Recompute fences to identify individual outlier wells
    df = oneline.copy()
    df['_eur'] = pd.to_numeric(df[eur_col], errors='coerce')
    if 'LateralLength' in df.columns:
        df['_lat'] = pd.to_numeric(df['LateralLength'], errors='coerce')
        df = df.dropna(subset=['_eur', '_lat'])
    else:
        df = df.dropna(subset=['_eur'])
    q1, q3      = df['_eur'].quantile(0.25), df['_eur'].quantile(0.75)
    iqr         = q3 - q1
    upper_f     = q3 + 1.5 * iqr
    lower_f     = q1 - 1.5 * iqr
    flagged_df  = df[(df['_eur'] > upper_f) | (df['_eur'] < lower_f)].copy()
    flagged_df['Outlier Type'] = flagged_df['_eur'].apply(
        lambda v: '⬆ High EUR' if v > upper_f else '⬇ Low EUR'
    )

    sel_key = f"{fluid}_tc_selection"
    all_ids = set(oneline['WellID'].astype(str).tolist())
    if sel_key not in st.session_state:
        st.session_state[sel_key] = all_ids.copy()

    if flagged_df.empty:
        st.success(f"No outliers detected for {fluid}.")
        return

    st.markdown(f"**{len(flagged_df)} outlier(s) detected.** Check wells to exclude from the Type Curve, then click **Update TC**.")

    show_cols = [c for c in ['WellName', 'API10', 'LateralLength', eur_col] if c in flagged_df.columns]
    display   = flagged_df[['WellID', 'Outlier Type'] + show_cols].reset_index(drop=True)
    current_excluded = all_ids - st.session_state[sel_key]
    display.insert(0, "Exclude", display['WellID'].astype(str).isin(current_excluded).tolist())

    col_cfg = {"Exclude": st.column_config.CheckboxColumn("Exclude", default=False)}
    for c in ['Outlier Type'] + show_cols:
        col_cfg[c] = st.column_config.TextColumn(c, disabled=True)

    edited = st.data_editor(
        display.drop(columns=['WellID']),
        column_config=col_cfg,
        hide_index=True,
        use_container_width=True,
        key=f"{fluid}_outlier_editor",
    )

    excl_mask    = edited['Exclude'].values.astype(bool)
    excl_well_ids = set(display.loc[excl_mask, 'WellID'].astype(str).tolist())
    n_excl        = int(excl_mask.sum())

    btn_col, reset_col = st.columns([2, 1])
    with btn_col:
        if st.button(
            f"Update {fluid} Type Curve — remove {n_excl} outlier(s)" if n_excl else f"Select wells above to exclude",
            disabled=(n_excl == 0),
            type="primary",
            key=f"{fluid}_apply_outlier",
        ):
            st.session_state[sel_key] = all_ids - excl_well_ids
            st.success(f"{fluid}: {n_excl} well(s) removed from Type Curve. Scroll down to see the updated TC.")
            st.rerun()
    with reset_col:
        if st.button(f"Reset {fluid} selection", key=f"{fluid}_reset_outlier"):
            st.session_state[sel_key] = all_ids.copy()
            st.success(f"{fluid} Type Curve restored to all {len(all_ids)} wells.")
            st.rerun()

_out_tabs = st.tabs(["Oil", "Gas", "Water"])
with _out_tabs[0]: _outlier_section("Oil",   "Oil_oneline",   "EUR (Mbbl)")
with _out_tabs[1]: _outlier_section("Gas",   "Gas_oneline",   "EUR (MMcf)")
with _out_tabs[2]: _outlier_section("Water", "Water_oneline", "EUR (Mbbl water)")

st.markdown("---")

# ================= Type Wells =================
st.header("Type Wells — Summary")
tw_tabs = st.tabs(["Oil", "Gas", "Water"])

def _render_tw(fluid, on_key, mo_key, eur_col):
    if on_key not in st.session_state:
        st.info(f"Run {fluid} first.")
        return

    oneline_full = st.session_state[on_key]

    # --- Analog selection table ---
    sel_key = f"{fluid}_tc_selection"

    # Build display columns available in the oneline
    show_cols = [c for c in [
        'WellName', 'PrimaryFormation', 'LateralLength', eur_col
    ] if c in oneline_full.columns]

    # Initialise selection state: set of all WellIDs (all selected)
    all_ids = set(oneline_full['WellID'].astype(str).tolist())
    if sel_key not in st.session_state:
        st.session_state[sel_key] = all_ids

    # Build editable frame: "Select" checkbox + display columns
    saved_ids = st.session_state[sel_key]
    editor_df = oneline_full[show_cols].copy().reset_index(drop=True)
    editor_df.insert(0, "Select", oneline_full['WellID'].astype(str).isin(saved_ids).tolist())

    col_cfg = {"Select": st.column_config.CheckboxColumn("Select", default=True)}
    for c in show_cols:
        col_cfg[c] = st.column_config.TextColumn(c, disabled=True)

    # ── Two-column layout: table left, plot right ──
    left_col, right_col = st.columns([1, 1])

    with left_col:
        st.caption(f"**Analog Wells** — check/uncheck to include or exclude from the Type Curve ({len(oneline_full)} wells total)")
        edited = st.data_editor(
            editor_df,
            column_config=col_cfg,
            use_container_width=True,
            hide_index=True,
            key=f"{fluid}_analog_editor",
        )

    # Persist selection as a set of WellIDs (robust across reruns/index shifts)
    selected_mask = edited["Select"].values.astype(bool)
    oneline_filtered = oneline_full[selected_mask].reset_index(drop=True)
    st.session_state[sel_key] = set(oneline_filtered['WellID'].astype(str).tolist())
    n_sel = int(selected_mask.sum())
    n_tot = len(oneline_full)

    # Compute EUR stats upfront so P50 curve can be scaled to match
    eurs = []
    eur_stats = {}
    if n_sel > 0 and eur_col in oneline_filtered.columns:
        eurs = pd.to_numeric(oneline_filtered[eur_col], errors='coerce').dropna().tolist()
        eur_stats = compute_eur_stats(eurs)

    with right_col:
        if mo_key in st.session_state and n_sel > 0:
            sel_ids = set(oneline_filtered['WellID'].astype(str))
            monthly_filtered = st.session_state[mo_key][
                st.session_state[mo_key]['WellID'].astype(str).isin(sel_ids)
            ]
            curves, lines = build_type_curves_and_lines(
                monthly_filtered, fluid.lower(), oneline=oneline_filtered
            )
            # Scale every TC curve so its cumulative volume exactly matches
            # the corresponding EUR percentile from the stats table.
            # Curves are in Bbl or Mcf; EUR stats are in Mbbl or MMcf → ×1000.
            if not curves.empty:
                curves = curves.copy()
                for _pct, _stat in [('P10', 'P10'), ('P50', 'P50'), ('P90', 'P90')]:
                    if (_pct in curves.columns
                            and eur_stats.get(_stat)
                            and np.isfinite(eur_stats[_stat])):
                        _csum = curves[_pct].sum()
                        if _csum > 0:
                            curves[_pct] = curves[_pct] * (eur_stats[_stat] * 1000 / _csum)
            st.session_state[f"{fluid}_tc_curves"] = curves
            st.session_state[f"{fluid}_tc_lines"] = lines
            st.pyplot(plot_type_curves(curves, lines, fluid.lower()))

    # ── EUR summary below ──
    st.caption(f"**{n_sel} / {n_tot} wells selected**")
    if n_sel > 0:
        st.dataframe(
            format_df_2dec(eur_summary_table(fluid, eur_stats, "Mbbl" if fluid != "Gas" else "MMcf", int(norm_len))),
            use_container_width=True,
        )
    else:
        st.warning("No wells selected — select at least one well to compute statistics.")

with tw_tabs[0]: _render_tw("Oil",   "Oil_oneline",   "Oil_monthly",   "EUR (Mbbl)")
with tw_tabs[1]: _render_tw("Gas",   "Gas_oneline",   "Gas_monthly",   "EUR (MMcf)")
with tw_tabs[2]: _render_tw("Water", "Water_oneline", "Water_monthly", "EUR (Mbbl water)")

# ── Multi-phase TC download ──
_tc_keys = ["Oil_tc_curves", "Gas_tc_curves", "Water_tc_curves"]
if any(k in st.session_state for k in _tc_keys):
    import io as _io
    _oil_c   = st.session_state.get("Oil_tc_curves",   pd.DataFrame())
    _gas_c   = st.session_state.get("Gas_tc_curves",   pd.DataFrame())
    _water_c = st.session_state.get("Water_tc_curves", pd.DataFrame())

    # Build combined Oil + Gas sheet (P50 columns adjacent, units converted)
    _og_parts = [pd.Series(range(1, 601), name="Month")]
    for _pct in ["P10", "P50", "P90"]:
        if not _oil_c.empty and _pct in _oil_c.columns:
            _og_parts.append((_oil_c[_pct] / 1000).round(4).rename(f"{_pct}_Oil_Mbbl_per_month"))
        if not _gas_c.empty and _pct in _gas_c.columns:
            _og_parts.append((_gas_c[_pct] / 1000).round(4).rename(f"{_pct}_Gas_MMcf_per_month"))
    _og_df = pd.concat(_og_parts, axis=1)

    # Water sheet
    _water_parts = [pd.Series(range(1, 601), name="Month")]
    for _pct in ["P10", "P50", "P90"]:
        if not _water_c.empty and _pct in _water_c.columns:
            _water_parts.append((_water_c[_pct] / 1000).round(4).rename(f"{_pct}_Water_Mbbl_per_month"))
    _water_df = pd.concat(_water_parts, axis=1)

    _xls_buf = _io.BytesIO()
    with pd.ExcelWriter(_xls_buf, engine="openpyxl") as _writer:
        _og_df.to_excel(_writer, sheet_name="Oil & Gas TC", index=False)
        if len(_water_parts) > 1:
            _water_df.to_excel(_writer, sheet_name="Water TC", index=False)
    st.download_button(
        label="⬇ Download All Phases TC Monthly Volumes (Excel)",
        data=_xls_buf.getvalue(),
        file_name="TC_Monthly_Volumes.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        key="all_tc_download",
    )

# ================= PDF EXPORT =================
from datetime import datetime


# ── PPTX data builder ──────────────────────────────────────────────────────────
def _build_pptx_data(tc_name: str) -> dict:
    """Extract live session-state data into the pptxgenjs data contract."""
    import json as _json
    from datetime import datetime as _dt

    D_LIM = 0.00417

    def _get_fluid(fluid_name, oneline_key, eur_col):
        if oneline_key not in st.session_state:
            return None
        oneline_full = st.session_state[oneline_key]
        sel_key = f"{fluid_name}_tc_selection"
        if sel_key in st.session_state and isinstance(st.session_state[sel_key], set):
            sel_ids = st.session_state[sel_key]
            oneline = oneline_full[oneline_full['WellID'].astype(str).isin(sel_ids)].reset_index(drop=True)
        else:
            oneline = oneline_full.reset_index(drop=True)

        n_tc    = len(oneline)
        n_total = len(oneline_full)
        is_gas  = fluid_name.lower() == 'gas'
        eur_unit      = "MMcf" if is_gas else "Mbbl"
        eur_per_ft_u  = "Mcf/ft" if is_gas else "Bbl/ft"
        qi_unit       = "Mcf/day" if is_gas else "bbl/day"

        eurs = pd.to_numeric(oneline[eur_col], errors='coerce').dropna().astype(float) if eur_col in oneline.columns else pd.Series([], dtype=float)
        eur_stats = compute_eur_stats(eurs.tolist()) if len(eurs) > 0 else {}
        _pk = lambda kw: next((k for k in eur_stats if kw in k.lower()), None)
        p10e  = float(eur_stats.get(_pk('p10'), 0) or 0)
        p50e  = float(eur_stats.get(_pk('p50') or _pk('median'), 0) or 0)
        p90e  = float(eur_stats.get(_pk('p90'), 0) or 0)
        mne   = float(eur_stats.get(_pk('mean'), 0) or 0)

        lat_vals = pd.to_numeric(oneline['LateralLength'], errors='coerce').dropna() if 'LateralLength' in oneline.columns else pd.Series([], dtype=float)
        avg_lat  = float(lat_vals.mean()) if len(lat_vals) > 0 else float(norm_len)
        def _epft(e): return (e * 1000) / avg_lat if avg_lat > 0 else 0.0

        b_vals = pd.to_numeric(oneline['b'], errors='coerce').dropna() if 'b' in oneline.columns else pd.Series([], dtype=float)
        bc = len(b_vals)

        # Group parameters (mirror build_type_curves_and_lines grouping)
        df_tc = oneline.copy()
        df_tc['_eur'] = pd.to_numeric(df_tc[eur_col], errors='coerce') if eur_col in df_tc.columns else 0
        df_tc['_qi']  = pd.to_numeric(df_tc.get('qi (per day)', pd.Series(dtype=float)), errors='coerce')
        df_tc['_b']   = pd.to_numeric(df_tc.get('b',            pd.Series(dtype=float)), errors='coerce')
        df_tc['_di']  = pd.to_numeric(df_tc.get('di (per month)', pd.Series(dtype=float)), errors='coerce')
        df_tc['_dec'] = pd.to_numeric(df_tc.get('First-Year Decline (%)', pd.Series(dtype=float)), errors='coerce')
        df_tc = df_tc.dropna(subset=['_eur','_qi','_b','_di']).sort_values('_eur', ascending=False).reset_index(drop=True)
        nn = len(df_tc)

        def _gs(grp):
            qi  = float(grp['_qi'].median()) if len(grp) > 0 else 0.0
            b_  = float(grp['_b'].median())  if len(grp) > 0 else 1.0
            di_ = float(grp['_di'].median()) if len(grp) > 0 else D_LIM
            dec = float(grp['_dec'].dropna().median()) / 100 if (len(grp) > 0 and grp['_dec'].notna().any()) else 0.0
            return qi, b_, di_, dec

        if nn >= 1:
            p50_lo = max(0, int(nn * 0.45))
            p50_hi = max(p50_lo + 1, int(nn * 0.55))
            grps = {
                'P10': df_tc.iloc[:max(1, int(nn * 0.1))],
                'P50': df_tc.iloc[p50_lo:p50_hi],
                'P90': df_tc.iloc[max(0, int(nn * 0.9)):],
            }
            p10qi, p10b, p10di, p10dec = _gs(grps['P10'])
            p50qi, p50b, p50di, p50dec = _gs(grps['P50'])
            p90qi, p90b, p90di, p90dec = _gs(grps['P90'])
        else:
            p10qi = p50qi = p90qi = 0.0
            p10b  = p50b  = p90b  = 1.0
            p10di = p50di = p90di = D_LIM
            p10dec = p50dec = p90dec = 0.0

        analog_rows = []
        for _, row in oneline.iterrows():
            api_v  = str(row.get('API10', row.get('WellID', '')))
            wn     = str(row.get('WellName', ''))
            lat_ft = float(row['LateralLength']) if 'LateralLength' in row and pd.notna(row['LateralLength']) else avg_lat
            qi_d   = float(row['qi (per day)'])        if 'qi (per day)'        in row and pd.notna(row['qi (per day)'])        else 0.0
            bv     = float(row['b'])                   if 'b'                   in row and pd.notna(row['b'])                   else 1.0
            di_mo  = float(row['di (per month)'])      if 'di (per month)'      in row and pd.notna(row['di (per month)'])      else D_LIM
            dec1   = float(row['First-Year Decline (%)']) if 'First-Year Decline (%)' in row and pd.notna(row['First-Year Decline (%)']) else 0.0
            eur_v  = float(row[eur_col])               if eur_col               in row and pd.notna(row[eur_col])               else 0.0
            epft_v = (eur_v * 1000) / lat_ft if lat_ft > 0 else 0.0
            analog_rows.append({
                'api': api_v, 'wellName': wn, 'latLenFt': lat_ft,
                'qi': qi_d, 'b': bv, 'diPerMo': di_mo,
                'decline1yrPct': dec1, 'eur': eur_v, 'eurPerFt': epft_v,
            })

        return {
            'fluid': fluid_name.lower(),
            'nTcWells': n_tc, 'nTotal': n_total,
            'meanEur': mne, 'p10Eur': p10e, 'p50Eur': p50e, 'p90Eur': p90e,
            'eurUnit': eur_unit, 'eurPerFtUnit': eur_per_ft_u,
            'meanEurPerFt': _epft(mne), 'p10EurPerFt': _epft(p10e),
            'p50EurPerFt': _epft(p50e), 'p90EurPerFt': _epft(p90e),
            'bMean':   float(b_vals.mean())             if bc > 0 else 0.0,
            'bMedian': float(b_vals.median())           if bc > 0 else 0.0,
            'bP10':    float(np.percentile(b_vals, 10)) if bc > 0 else 0.0,
            'bP90':    float(np.percentile(b_vals, 90)) if bc > 0 else 0.0,
            'bCount': int(bc),
            'bValues': [float(v) for v in b_vals.tolist()],
            'p90Qi': p90qi, 'p50Qi': p50qi, 'p10Qi': p10qi, 'qiUnit': qi_unit,
            'p90B':  p90b,  'p50B':  p50b,  'p10B':  p10b,
            'p90Di': p90di, 'p50Di': p50di, 'p10Di': p10di,
            'p90Decline1yr': p90dec, 'p50Decline1yr': p50dec, 'p10Decline1yr': p10dec,
            'terminalDi': D_LIM,
            'analogRows': analog_rows,
        }

    # Well meta — derive from best available oneline
    primary_key = next((k for k in ['Oil_oneline','Gas_oneline','Water_oneline'] if k in st.session_state), None)
    prim_fl     = primary_key.split('_')[0] if primary_key else 'Oil'
    prim_ol     = st.session_state.get(primary_key, pd.DataFrame())
    sel_key     = f"{prim_fl}_tc_selection"
    if sel_key in st.session_state and isinstance(st.session_state[sel_key], set):
        prim_tc = prim_ol[prim_ol['WellID'].astype(str).isin(st.session_state[sel_key])]
    else:
        prim_tc = prim_ol

    def _mode(col): return str(prim_tc[col].mode().iloc[0]) if (col in prim_tc.columns and not prim_tc[col].mode().empty) else ''
    county    = _mode('County')
    state_val = _mode('State')
    reservoir = _mode('PrimaryFormation')
    lat_vals  = pd.to_numeric(prim_tc['LateralLength'], errors='coerce').dropna() if 'LateralLength' in prim_tc.columns else pd.Series([], dtype=float)
    avg_lat   = float(lat_vals.mean()) if len(lat_vals) > 0 else float(norm_len)

    comp_dates  = pd.to_datetime(prim_tc['CompletionDate'], errors='coerce').dropna() if 'CompletionDate' in prim_tc.columns else pd.Series([], dtype='datetime64[ns]')
    comp_vintage = (f"{comp_dates.min().year} – {comp_dates.max().year}" if len(comp_dates) > 0 else '')

    unique_counties = list(prim_tc['County'].dropna().unique()) if 'County' in prim_tc.columns else [county]
    ll_range = (f"{int(lat_vals.min()):,} – {int(lat_vals.max()):,} ft" if len(lat_vals) > 0 else '')

    now = _dt.now()
    data = {
        'wellMeta': {
            'wellName':       tc_name or 'TC Report',
            'county':         county,
            'state':          state_val,
            'operator':       'SE',
            'reservoir':      reservoir,
            'afeLateralLenFt': float(norm_len),
            'tvdFt':          0.0,
            'reportDate':     now.isoformat(),
            'generatedDate':  now.strftime('%Y-%m-%d'),
            'normLengthFt':   float(norm_len),
            'bFactorRange':   f"{b_low:.3f} - {b_high:.3f}",
        },
        'analogCriteria': {
            'nAnalogs':        len(prim_tc),
            'counties':        ', '.join(str(c) for c in unique_counties),
            'reservoir':       reservoir,
            'llRangeFt':       ll_range,
            'llAvgFt':         avg_lat,
            'compVintage':     comp_vintage,
            'avgDistanceMi':   'N/A',
        },
    }
    _blank = lambda fl, eu, epfu, qu: {
        'fluid': fl, 'nTcWells': 0, 'nTotal': 0,
        'meanEur': 0.0, 'p10Eur': 0.0, 'p50Eur': 0.0, 'p90Eur': 0.0,
        'eurUnit': eu, 'eurPerFtUnit': epfu,
        'meanEurPerFt': 0.0, 'p10EurPerFt': 0.0, 'p50EurPerFt': 0.0, 'p90EurPerFt': 0.0,
        'bMean': 0.0, 'bMedian': 0.0, 'bP10': 0.0, 'bP90': 0.0, 'bCount': 0, 'bValues': [],
        'p90Qi': 0.0, 'p50Qi': 0.0, 'p10Qi': 0.0, 'qiUnit': qu,
        'p90B': 1.0, 'p50B': 1.0, 'p10B': 1.0,
        'p90Di': 0.00417, 'p50Di': 0.00417, 'p10Di': 0.00417,
        'p90Decline1yr': 0.0, 'p50Decline1yr': 0.0, 'p10Decline1yr': 0.0,
        'terminalDi': 0.00417, 'analogRows': [],
    }
    data['oil']   = _get_fluid('Oil',   'Oil_oneline',   'EUR (Mbbl)')       or _blank('oil',   'Mbbl',  'Bbl/ft', 'bbl/day')
    data['gas']   = _get_fluid('Gas',   'Gas_oneline',   'EUR (MMcf)')       or _blank('gas',   'MMcf',  'Mcf/ft', 'Mcf/day')
    data['water'] = _get_fluid('Water', 'Water_oneline', 'EUR (Mbbl water)') or _blank('water', 'Mbbl',  'Bbl/ft', 'bbl/day')

    # Generate the exact plots shown in the Streamlit interface as PNG bytes
    def _fig_bytes(fig):
        buf = BytesIO()
        fig.savefig(buf, format='png', dpi=180, bbox_inches='tight')
        buf.seek(0); plt.close(fig)
        return buf.read()

    def _attach_figs(fluid_name, oneline_key, monthly_key, eur_col, fd):
        if oneline_key not in st.session_state or fd['nTcWells'] == 0:
            return
        oneline_full = st.session_state[oneline_key]
        sel_key = f"{fluid_name}_tc_selection"
        if sel_key in st.session_state and isinstance(st.session_state[sel_key], set):
            oneline = oneline_full[oneline_full['WellID'].astype(str).isin(st.session_state[sel_key])]
        else:
            oneline = oneline_full

        # B-factor histogram (same as bfactor_analytics_figures but keep figure)
        b_vals = pd.to_numeric(oneline.get('b'), errors='coerce').dropna()
        if len(b_vals) > 0:
            color = _phase_color(fluid_name)
            P50 = float(np.median(b_vals))
            fig, ax = plt.subplots(figsize=(6.6, 3.4))
            ax.hist(b_vals.values, bins=30, color=color, alpha=0.85, edgecolor='none')
            ax.axvline(P50, color='black', linestyle='--', label=f"Median={P50:.2f}")
            ax.legend(); ax.set_xlabel('b-factor'); ax.set_ylabel('Count')
            plt.tight_layout()
            fd['bfactor_hist_bytes'] = _fig_bytes(fig)

        # Type curve — use cached session-state curves so it matches the UI exactly
        curves = st.session_state.get(f"{fluid_name}_tc_curves", pd.DataFrame())
        lines  = st.session_state.get(f"{fluid_name}_tc_lines",  [])
        if not curves.empty:
            fig = plot_type_curves(curves, lines, fluid_name.lower())
            fd['typecurve_bytes'] = _fig_bytes(fig)

        # Probit plot
        if eur_col in oneline.columns:
            eurs = pd.to_numeric(oneline[eur_col], errors='coerce').dropna().astype(float).tolist()
            if eurs:
                unit = "Mbbl" if fluid_name != "Gas" else "MMcf"
                fig = probit_plot(eurs, unit, f"{fluid_name} EUR Distribution",
                                  _phase_color(fluid_name), norm_len=int(norm_len))
                fd['probit_bytes'] = _fig_bytes(fig)

    _attach_figs('Oil',   'Oil_oneline',   'Oil_monthly',   'EUR (Mbbl)',       data['oil'])
    _attach_figs('Gas',   'Gas_oneline',   'Gas_monthly',   'EUR (MMcf)',        data['gas'])
    _attach_figs('Water', 'Water_oneline', 'Water_monthly', 'EUR (Mbbl water)', data['water'])
    return data


def generate_tc_pptx(pptx_data: dict, output_dir: str) -> str | None:
    """Generate a TC PowerPoint report using python-pptx (no Node.js required)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        st.warning("python-pptx not installed. Run: pip install python-pptx")
        return None

    import math as _math
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as _plt
    from io import BytesIO as _BytesIO
    from datetime import datetime as _dt

    # ── Brand constants ────────────────────────────────────────────────────────
    SE_RED    = RGBColor(0x95, 0x37, 0x35)
    NAVY      = RGBColor(0x1A, 0x27, 0x44)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    ROW_ALT   = RGBColor(0xF5, 0xED, 0xED)
    BODY_GRAY = RGBColor(0x40, 0x40, 0x40)
    MARK_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
    FLUID_HEX = {'oil': '#2D6E2D', 'gas': '#CC3333', 'water': '#2244AA'}
    D_LIM     = 0.00417
    _RECT     = 1   # MSO_AUTO_SHAPE_TYPE.RECTANGLE
    _OVAL     = 9   # MSO_AUTO_SHAPE_TYPE.OVAL
    _ALIGN    = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    _blank = prs.slide_layouts[6]

    # ── Drawing primitives ─────────────────────────────────────────────────────
    def _rect(sl, x, y, w, h, fill, line=None, lw=0.5):
        sp = sl.shapes.add_shape(_RECT, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line:
            sp.line.color.rgb = line; sp.line.width = Pt(lw)
        else:
            sp.line.fill.background()
        return sp

    def _oval(sl, x, y, w, h, fill, line, lw=1.5):
        sp = sl.shapes.add_shape(_OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
        return sp

    def _txt(sl, text, x, y, w, h, size=10, bold=False, italic=False,
             color=None, face="Garamond", align='left', wrap=True):
        tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.alignment = _ALIGN.get(align, PP_ALIGN.LEFT)
        r = p.add_run(); r.text = str(text)
        r.font.size = Pt(size); r.font.bold = bool(bold); r.font.italic = bool(italic)
        r.font.name = face
        if color: r.font.color.rgb = color
        return tb

    def _chrome(sl, pg):
        _rect(sl, 0, 5.3, 10, 0.325, SE_RED)
        _oval(sl, 0.15, 5.05, 0.5, 0.5, NAVY, SE_RED)
        _txt(sl, "SE",  0.15, 5.12, 0.5,  0.35, size=9,  bold=True, color=WHITE, face="Garamond", align='center')
        _txt(sl, str(pg), 9.5, 5.3,  0.4,  0.32, size=10, color=WHITE, align='right')
        _txt(sl, "CONFIDENTIAL", 0.0, 0.5, 0.5, 4.5, size=7, color=MARK_GRAY, align='center')
        _txt(sl, "CONFIDENTIAL", 9.5, 0.5, 0.5, 4.5, size=7, color=MARK_GRAY, align='center')

    def _slide_title(sl, title):
        _txt(sl, title, 0.4, 0.18, 9.2, 0.52, size=24, bold=True, color=SE_RED, face="Garamond")
        _rect(sl, 0.4, 0.73, 9.2, 0.03, SE_RED)

    def _kv_box(sl, header, rows, x, y, w, rh=0.30):
        _rect(sl, x, y, w, rh, SE_RED)
        _txt(sl, header, x, y, w, rh, size=10, bold=True, color=WHITE)
        for i, (lbl, val) in enumerate(rows):
            ry = y + rh * (i + 1)
            bg = WHITE if i % 2 == 0 else ROW_ALT
            _rect(sl, x, ry, w, rh, bg, RGBColor(0xDD, 0xDD, 0xDD))
            _txt(sl, lbl,       x,          ry, w * 0.5, rh, size=9.5, color=BODY_GRAY)
            _txt(sl, str(val),  x + w * 0.5, ry, w * 0.5, rh, size=9.5, color=BODY_GRAY, align='center')

    def _table(sl, headers, rows, x, y, cws, rh=0.28):
        cx = x
        for i, hdr in enumerate(headers):
            _rect(sl, cx, y, cws[i], rh, SE_RED)
            _txt(sl, hdr, cx, y, cws[i], rh, size=9, bold=True, color=WHITE, align='center')
            cx += cws[i]
        for ri, row in enumerate(rows):
            ry = y + rh * (ri + 1)
            bg = WHITE if ri % 2 == 0 else ROW_ALT
            cx2 = x
            for ci, cell in enumerate(row):
                _rect(sl, cx2, ry, cws[ci], rh, bg, RGBColor(0xDD, 0xDD, 0xDD))
                _txt(sl, str(cell), cx2, ry, cws[ci], rh, size=9, color=BODY_GRAY, align='center')
                cx2 += cws[ci]

    def _embed(sl, fig, x, y, w, h):
        buf = _BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='#F8F8F8')
        buf.seek(0); _plt.close(fig)
        sl.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

    # ── Format helpers ─────────────────────────────────────────────────────────
    def _ff(n):       return f"{int(round(n)):,} ft"
    def _fe(n, u):    return f"{n:.1f} {u}"
    def _fn(n):       return f"{int(round(n)):,}" if n >= 1000 else f"{n:.2f}"
    def _fp(v):       return f"{v * 100:.1f}%"
    def _cap(s):      return s[0].upper() + s[1:] if s else s
    def _fd(iso):
        try:    return _dt.fromisoformat(iso).strftime("%B %-d, %Y")
        except: return str(iso)

    def _arps(qi_d, b, di, term_di, months=600):
        qi = qi_d * 30.4
        t_sw = ((di - term_di) / (b * di * term_di)) if (b > 0 and di > term_di) else float('inf')
        q_sw = qi / (1 + b * di * t_sw) ** (1 / b) if t_sw < float('inf') else qi
        q = []
        for t in range(1, months + 1):
            if t >= t_sw and t_sw < float('inf'):
                q.append(q_sw * _math.exp(-term_di * (t - t_sw)))
            else:
                q.append(qi / (1 + b * di * t) ** (1 / b))
        return q

    wm  = pptx_data['wellMeta']
    ac  = pptx_data['analogCriteria']
    oil = pptx_data['oil']
    gas = pptx_data['gas']
    wtr = pptx_data['water']

    # ── Slide 1: Title ─────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(_blank)
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = NAVY
    _oval(sl, 0.3, 0.2, 1.1, 1.1, NAVY, SE_RED, lw=3)
    _txt(sl, "SE", 0.3, 0.42, 1.1, 0.5, size=22, bold=True, color=WHITE, face="Garamond", align='center')
    _txt(sl, wm['wellName'], 0.4, 1.8, 7, 0.85, size=36, bold=True, color=WHITE, face="Garamond")
    _txt(sl, f"TC Report – {_fd(wm['reportDate'])}", 0.4, 2.7, 7, 0.6,
         size=24, italic=True, color=RGBColor(0xCC, 0xCC, 0xCC), face="Garamond")
    _rect(sl, 0, 5.3, 10, 0.325, SE_RED)

    # ── Slide 2: Disclaimer ────────────────────────────────────────────────────
    sl = prs.slides.add_slide(_blank)
    _slide_title(sl, "Disclaimer"); _chrome(sl, 2)
    _txt(sl, (
        'The information contained in this confidential presentation (this “Presentation”) is provided for '
        'informational and discussion purposes only and is not, and may not be relied on in any manner as, '
        'legal, tax or investment advice or as an offer to sell or a solicitation of an offer to buy an interest '
        'in any security. The information contained in this Presentation must be kept strictly confidential and '
        'may not be reproduced or redistributed in any format without the approval of Schaper International '
        'Petroleum Consulting, LLC (“SIPC”). In considering any performance data contained in this Presentation, '
        'you should bear in mind that past or targeted performance is not indicative of future results, and there '
        'can be no assurance that any investment will achieve comparable results or that target returns will be met. '
        'In addition, there can be no assurance that any investment will achieve or be realized at the valuations '
        'shown, as actual realized returns will depend on, among other factors, future operating results, the value '
        'of assets and market conditions at the time of disposition, any related transaction costs and the timing '
        'and manner of sale, all of which may differ from the assumptions on which the valuations contained in this '
        'Presentation are based. Nothing contained in this Presentation should be deemed to be a prediction or '
        'projection of future performance of any investment. Investors should make their own investigations and '
        'evaluations of a potential investment and the information contained in this Presentation. Except where '
        'otherwise indicated in this Presentation, the information provided in this Presentation is based on matters '
        'as they exist as of the date of preparation and not as of any future date, and will not be updated or '
        'otherwise revised to reflect information that subsequently becomes available, or circumstances existing or '
        'changes occurring after the date hereof.'
    ), 0.4, 0.9, 9.2, 4.3, size=10.5, color=RGBColor(0x33, 0x33, 0x33), wrap=True)

    # ── Slide 3: Analog Selection ──────────────────────────────────────────────
    sl = prs.slides.add_slide(_blank)
    _slide_title(sl, "Analog Selection Methodology"); _chrome(sl, 3)
    _kv_box(sl, "Subject Location", [
        ["Well",      wm['wellName']],
        ["County",    wm.get('county', '')],
        ["Reservoir", wm.get('reservoir', '')],
        ["AFE LL",    _ff(wm['afeLateralLenFt'])],
    ], 0.4, 0.85, 3.5)
    _kv_box(sl, "Analog Pool Criteria", [
        ["# of Analogs",            str(ac.get('nAnalogs', ''))],
        ["County",                  str(ac.get('counties', ''))],
        ["Reservoir",               str(ac.get('reservoir', ''))],
        ["LL Range",                str(ac.get('llRangeFt', ''))],
        ["LL Avg",                  _ff(ac['llAvgFt']) if ac.get('llAvgFt') else ''],
        ["Completion Vintage",      str(ac.get('compVintage', ''))],
        ["Avg Distance to Subject", str(ac.get('avgDistanceMi', 'N/A'))],
    ], 0.4, 2.25, 3.5)
    _rect(sl, 4.2, 0.85, 5.4, 4.2, RGBColor(0xCC, 0xCC, 0xCC), RGBColor(0x99, 0x99, 0x99))
    _txt(sl, "[ Insert analog location map here ]", 4.2, 2.5, 5.4, 0.6,
         size=11, color=RGBColor(0x66, 0x66, 0x66), italic=True, align='center')

    # ── Slide 4: Summary ───────────────────────────────────────────────────────
    sl = prs.slides.add_slide(_blank)
    _slide_title(sl, f"Summary – {wm['wellName']} TC"); _chrome(sl, 4)
    _txt(sl, "Analysis Parameters", 0.4, 0.85, 9, 0.3, size=12, bold=True, color=BODY_GRAY)
    _kv_box(sl, "Parameter", [
        ["Type Curve Name",      f"{wm['wellName']} TC"],
        ["Generated",            wm['generatedDate']],
        ["Normalization Length", _ff(wm['normLengthFt'])],
        ["B-factor Range",       wm.get('bFactorRange', '')],
    ], 0.4, 1.18, 4.5)
    _txt(sl, "EUR Summary", 0.4, 2.7, 9, 0.3, size=12, bold=True, color=BODY_GRAY)

    def _eur_row(fd):
        return [_cap(fd['fluid']),
                f"{fd['nTcWells']} / {fd['nTotal']}",
                _fe(fd['p90Eur'], fd['eurUnit']),
                _fe(fd['p50Eur'], fd['eurUnit']),
                _fe(fd['p10Eur'], fd['eurUnit']),
                _fe(fd['meanEur'], fd['eurUnit'])]

    _table(sl, ["Fluid", "TC Wells", "P90 EUR", "P50 EUR", "P10 EUR", "Mean EUR"],
           [_eur_row(oil), _eur_row(gas), _eur_row(wtr)],
           0.4, 3.02, [1.1, 0.9, 1.65, 1.65, 1.65, 1.65])

    # ── Slides 5/8/11: Fluid stats + parameters ────────────────────────────────
    for fd, pg in [(oil, 5), (gas, 8), (wtr, 11)]:
        if not fd or fd['nTcWells'] == 0:
            continue
        sl = prs.slides.add_slide(_blank)
        _slide_title(sl, f"{_cap(fd['fluid'])} Analysis"); _chrome(sl, pg)
        _txt(sl, f"{_cap(fd['fluid'])} EUR Statistics ({fd['eurUnit']})",
             0.4, 0.85, 3.5, 0.3, size=11, bold=True, color=BODY_GRAY)
        _kv_box(sl, "Metric", [
            ["Wells in TC",  f"{fd['nTcWells']} / {fd['nTotal']}"],
            ["Mean",         f"{fd['meanEur']:.2f}"],
            ["Median (P50)", f"{fd['p50Eur']:.2f}"],
            ["P10",          f"{fd['p10Eur']:.2f}"],
            ["P90",          f"{fd['p90Eur']:.2f}"],
        ], 0.4, 1.17, 3.5)
        _txt(sl, "B-Factor Statistics", 0.4, 3.0, 3.5, 0.3, size=11, bold=True, color=BODY_GRAY)
        _kv_box(sl, "Metric", [
            ["Count",    str(fd['bCount'])],
            ["Mean",     f"{fd['bMean']:.3f}"],
            ["Median",   f"{fd['bMedian']:.3f}"],
            ["P10/P90",  f"{fd['bP10']:.3f} / {fd['bP90']:.3f}"],
        ], 0.4, 3.3, 3.5)
        _txt(sl, "Type Curve Parameters", 4.1, 0.85, 5.6, 0.3, size=11, bold=True, color=BODY_GRAY)
        _table(sl, ["Parameter", "P90 (Low)", "P50 (Mid)", "P10 (High)"], [
            [f"qi ({fd['qiUnit']})", _fn(fd['p90Qi']),  _fn(fd['p50Qi']),  _fn(fd['p10Qi'])],
            ["b-factor",             f"{fd['p90B']:.3f}", f"{fd['p50B']:.3f}", f"{fd['p10B']:.3f}"],
            ["Di (per month)",       f"{fd['p90Di']:.4f}", f"{fd['p50Di']:.4f}", f"{fd['p10Di']:.4f}"],
            ["1st-Year Decline",     _fp(fd['p90Decline1yr']), _fp(fd['p50Decline1yr']), _fp(fd['p10Decline1yr'])],
            ["Terminal Di",          f"{fd['terminalDi']:.5f}/mo", f"{fd['terminalDi']:.5f}/mo", f"{fd['terminalDi']:.5f}/mo"],
            [f"EUR ({fd['eurUnit']})", f"{fd['p90Eur']:.2f}", f"{fd['p50Eur']:.2f}", f"{fd['p10Eur']:.2f}"],
        ], 4.1, 1.17, [2.2, 1.1, 1.1, 1.2])
        _txt(sl, f"{_cap(fd['fluid'])} EURs", 4.1, 3.4, 5.6, 0.3, size=11, bold=True, color=BODY_GRAY)
        _table(sl, ["Percentile", f"EUR ({fd['eurUnit']})", f"EUR ({fd['eurPerFtUnit']})"], [
            ["P90",  f"{fd['p90Eur']:.2f}",  f"{fd['p90EurPerFt']:.2f}"],
            ["P50",  f"{fd['p50Eur']:.2f}",  f"{fd['p50EurPerFt']:.2f}"],
            ["P10",  f"{fd['p10Eur']:.2f}",  f"{fd['p10EurPerFt']:.2f}"],
            ["Mean", f"{fd['meanEur']:.2f}", f"{fd['meanEurPerFt']:.2f}"],
        ], 4.1, 3.72, [1.5, 2.05, 2.05])

    # ── Slides 6/9/12: Charts — paste exact Streamlit interface plots as images ──
    def _embed_bytes(sl, data_bytes, x, y, w, h):
        sl.shapes.add_picture(_BytesIO(data_bytes), Inches(x), Inches(y), Inches(w), Inches(h))

    for fd, pg in [(oil, 6), (gas, 9), (wtr, 12)]:
        if not fd or fd['nTcWells'] == 0:
            continue
        sl = prs.slides.add_slide(_blank)
        _slide_title(sl, f"{_cap(fd['fluid'])} Analysis – Charts"); _chrome(sl, pg)

        has_hist = 'bfactor_hist_bytes' in fd
        has_tc   = 'typecurve_bytes' in fd
        has_prob = 'probit_bytes' in fd

        if has_hist and has_tc:
            # Top row: b-factor histogram (left) + type curve (right)
            _txt(sl, "B-Factor Distribution", 0.4, 0.82, 4.5, 0.25, size=10, bold=True, color=BODY_GRAY)
            _embed_bytes(sl, fd['bfactor_hist_bytes'], 0.4, 1.08, 4.5, 2.1)
            _txt(sl, "Type Curve", 5.1, 0.82, 4.5, 0.25, size=10, bold=True, color=BODY_GRAY)
            _embed_bytes(sl, fd['typecurve_bytes'], 5.1, 1.08, 4.5, 2.1)
        elif has_hist:
            _txt(sl, "B-Factor Distribution", 0.4, 0.82, 9.2, 0.25, size=10, bold=True, color=BODY_GRAY)
            _embed_bytes(sl, fd['bfactor_hist_bytes'], 0.4, 1.08, 9.2, 2.1)
        elif has_tc:
            _txt(sl, "Type Curve", 0.4, 0.82, 9.2, 0.25, size=10, bold=True, color=BODY_GRAY)
            _embed_bytes(sl, fd['typecurve_bytes'], 0.4, 1.08, 9.2, 2.1)

        if has_prob:
            _txt(sl, "EUR Probit", 0.4, 3.25, 9.2, 0.25, size=10, bold=True, color=BODY_GRAY)
            _embed_bytes(sl, fd['probit_bytes'], 0.4, 3.5, 9.2, 1.8)

    # ── Slides 7/10/13: Analog well tables ─────────────────────────────────────
    for fd, pg in [(oil, 7), (gas, 10), (wtr, 13)]:
        if not fd or fd['nTcWells'] == 0:
            continue
        sl = prs.slides.add_slide(_blank)
        _slide_title(sl, f"{_cap(fd['fluid'])} Analysis – Analogs"); _chrome(sl, pg)
        rows_data = fd.get('analogRows', [])
        n = len(rows_data)
        rh = 0.225 if n > 15 else (0.255 if n > 10 else 0.31)
        _table(sl,
               ["API/UWI", "Well Name", "Lat Len (ft)",
                f"qi ({fd['qiUnit']})", "b", "Di (/mo)", "1yr Dec (%)", f"EUR ({fd['eurUnit']})"],
               [[str(r['api']), str(r['wellName']), f"{int(round(r['latLenFt'])):,}",
                 _fn(r['qi']), f"{r['b']:.4f}", f"{r['diPerMo']:.4f}",
                 f"{r['decline1yrPct']:.2f}", f"{r['eur']:.2f}"] for r in rows_data],
               0.4, 0.88, [1.35, 1.85, 1.0, 0.85, 0.8, 0.75, 0.9, 0.7], rh)

    # ── Save ───────────────────────────────────────────────────────────────────
    try:
        safe = "".join(c if c.isalnum() or c in ' _-' else '_' for c in wm['wellName'])
        out_path = os.path.join(output_dir, f"SE_{safe or 'TC'}_TC_Report_{wm['generatedDate']}.pptx")
        prs.save(out_path)
        return out_path
    except Exception as e:
        st.warning(f"PPTX save failed: {e}")
        return None

def generate_comprehensive_pdf(tc_name: str = ""):
    """Generate a comprehensive PDF report with all sections."""
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    doc = SimpleDocTemplate(tmp.name, pagesize=letter, topMargin=0.5*inch, bottomMargin=0.5*inch)
    story = []
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle('CustomTitle', parent=styles['Title'], fontSize=20, textColor=colors.HexColor('#781F28'))
    heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading1'], fontSize=14, textColor=colors.HexColor('#1976D2'))
    tc_name_style = ParagraphStyle('TCName', parent=styles['Normal'], fontSize=16,
                                   textColor=colors.HexColor('#1976D2'), fontName='Helvetica-Bold')

    # Title page
    story.append(Paragraph("SE Oil & Gas Autoforecasting Report", title_style))
    story.append(Spacer(1, 12))
    if tc_name:
        name_data = [[Paragraph(f"Type Curve Name:&nbsp;&nbsp;<font color='#1976D2'>{tc_name}</font>", tc_name_style)]]
        name_box = Table(name_data, colWidths=[6.5*inch])
        name_box.setStyle(TableStyle([
            ('BACKGROUND',    (0, 0), (-1, -1), colors.HexColor('#EAF4FB')),
            ('TOPPADDING',    (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 10),
            ('LEFTPADDING',   (0, 0), (-1, -1), 14),
            ('LINEBEFORE',    (0, 0), (0, -1),  4, colors.HexColor('#1976D2')),
        ]))
        story.append(name_box)
        story.append(Spacer(1, 10))
    story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
    story.append(Spacer(1, 24))
    
    # Parameters section
    story.append(Paragraph("Analysis Parameters", heading_style))
    story.append(Spacer(1, 12))
    params_data = [
        ["Parameter", "Value"],
        ["Normalization Length", f"{norm_len} ft"],
        ["Apply Normalization", "Yes" if use_norm else "No"],
        ["B-factor Range", f"{b_low:.3f} - {b_high:.3f}"],
        ["Geographic Bin Decimals", str(bin_decimals)]
    ]
    params_table = Table(params_data, colWidths=[3*inch, 2*inch])
    params_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 10),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(params_table)
    story.append(Spacer(1, 20))

    # ── Executive Summary (right after parameters) ──
    story.append(Paragraph("Executive Summary", heading_style))
    story.append(Spacer(1, 8))

    exec_data = [["Fluid", "TC Wells", "P90 EUR", "P50 EUR", "P10 EUR", "Mean EUR"]]
    for _fl, _key, _ecol in [("Oil", "Oil_oneline", "EUR (Mbbl)"),
                               ("Gas", "Gas_oneline", "EUR (MMcf)"),
                               ("Water", "Water_oneline", "EUR (Mbbl water)")]:
        if _key not in st.session_state:
            continue
        _oneline_full = st.session_state[_key]
        _sel_key = f"{_fl}_tc_selection"
        if _sel_key in st.session_state and isinstance(st.session_state[_sel_key], set):
            _sel_ids = st.session_state[_sel_key]
            _oneline_tc = _oneline_full[_oneline_full['WellID'].astype(str).isin(_sel_ids)]
        else:
            _oneline_tc = _oneline_full
        _tc_count = len(_oneline_tc)
        if _ecol in _oneline_tc.columns:
            _eurs = pd.to_numeric(_oneline_tc[_ecol], errors="coerce").dropna().astype(float).tolist()
            if _eurs:
                _st = compute_eur_stats(_eurs)
                _unit = "MMcf" if _fl == "Gas" else "Mbbl"
                _p90k = next((k for k in _st if 'p90' in k.lower()), None)
                _p50k = next((k for k in _st if 'p50' in k.lower() or 'median' in k.lower()), None)
                _p10k = next((k for k in _st if 'p10' in k.lower()), None)
                _mnk  = next((k for k in _st if 'mean' in k.lower()), None)
                exec_data.append([
                    _fl, str(_tc_count),
                    f"{_st[_p90k]:.1f} {_unit}" if _p90k else "N/A",
                    f"{_st[_p50k]:.1f} {_unit}" if _p50k else "N/A",
                    f"{_st[_p10k]:.1f} {_unit}" if _p10k else "N/A",
                    f"{_st[_mnk]:.1f} {_unit}" if _mnk else "N/A",
                ])
    if len(exec_data) > 1:
        exec_tbl = Table(exec_data, colWidths=[0.8*inch, 0.8*inch, 1.3*inch, 1.3*inch, 1.3*inch, 1.3*inch])
        exec_tbl.setStyle(TableStyle([
            ('BACKGROUND',    (0, 0), (-1, 0),  colors.HexColor('#781F28')),
            ('TEXTCOLOR',     (0, 0), (-1, 0),  colors.whitesmoke),
            ('ALIGN',         (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME',      (0, 0), (-1, 0),  'Helvetica-Bold'),
            ('FONTSIZE',      (0, 0), (-1, -1), 9),
            ('BOTTOMPADDING', (0, 0), (-1, 0),  10),
            ('TOPPADDING',    (0, 0), (-1, -1), 5),
            ('GRID',          (0, 0), (-1, -1), 0.5, colors.black),
            ('ROWBACKGROUNDS',(0, 1), (-1, -1), [colors.white, colors.HexColor('#F9ECEC')]),
        ]))
        story.append(exec_tbl)
    story.append(PageBreak())

    # Function to add fluid section
    def add_fluid_section(fluid_name, oneline_key, monthly_key, eur_col):
        if oneline_key not in st.session_state:
            return

        story.append(Paragraph(f"{fluid_name} Analysis", heading_style))
        story.append(Spacer(1, 12))

        oneline_full = st.session_state[oneline_key]

        # Apply same well selection as shown in the TC section (stored as a set of WellIDs)
        sel_key = f"{fluid_name}_tc_selection"
        if sel_key in st.session_state and isinstance(st.session_state[sel_key], set):
            sel_ids = st.session_state[sel_key]
            oneline = oneline_full[oneline_full['WellID'].astype(str).isin(sel_ids)].reset_index(drop=True)
        else:
            oneline = oneline_full
            sel_ids = set(oneline_full['WellID'].astype(str).tolist())

        # Filter monthly to the same selected wells
        monthly_full = st.session_state.get(monthly_key)
        if monthly_full is not None:
            monthly_sel = monthly_full[monthly_full['WellID'].astype(str).isin(sel_ids)]
        else:
            monthly_sel = None

        # Well count
        well_count = len(oneline)
        story.append(Paragraph(f"Wells in Type Curve: {well_count} / {len(oneline_full)}", styles['Normal']))
        story.append(Spacer(1, 12))
        
        # EUR Statistics
        if eur_col in oneline.columns:
            eurs = pd.to_numeric(oneline[eur_col], errors="coerce").dropna().astype(float).tolist()

            if eurs:
                stats = compute_eur_stats(eurs)
                unit = "Mbbl" if fluid_name != "Gas" else "MMcf"
                
                story.append(Paragraph(f"EUR Statistics ({unit})", styles['Heading2']))
                story.append(Spacer(1, 6))
                
                # Build EUR data table with available stats
                eur_data = [["Metric", "Value"]]
                
                # Handle different possible key formats
                count_key = next((k for k in stats.keys() if 'count' in k.lower()), None)
                if count_key:
                    eur_data.append(["Count", f"{stats[count_key]:.0f}"])
                
                mean_key = next((k for k in stats.keys() if 'mean' in k.lower()), None)
                if mean_key:
                    eur_data.append(["Mean", f"{stats[mean_key]:.2f}"])
                
                # Try P50 or median
                p50_key = next((k for k in stats.keys() if 'p50' in k.lower() or 'median' in k.lower()), None)
                if p50_key:
                    eur_data.append(["Median (P50)", f"{stats[p50_key]:.2f}"])
                
                p10_key = next((k for k in stats.keys() if 'p10' in k.lower()), None)
                if p10_key:
                    eur_data.append(["P10", f"{stats[p10_key]:.2f}"])
                
                p90_key = next((k for k in stats.keys() if 'p90' in k.lower()), None)
                if p90_key:
                    eur_data.append(["P90", f"{stats[p90_key]:.2f}"])
                
                min_key = next((k for k in stats.keys() if 'min' in k.lower()), None)
                if min_key:
                    eur_data.append(["Min", f"{stats[min_key]:.2f}"])
                
                max_key = next((k for k in stats.keys() if 'max' in k.lower()), None)
                if max_key:
                    eur_data.append(["Max", f"{stats[max_key]:.2f}"])
                
                eur_table = Table(eur_data, colWidths=[2.5*inch, 2*inch])
                eur_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), _phase_color_rgb(fluid_name)),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(eur_table)
                story.append(Spacer(1, 18))
        
        # B-factor statistics
        if 'b' in oneline.columns:
            b_vals = pd.to_numeric(oneline['b'], errors='coerce').dropna()
            if not b_vals.empty:
                story.append(Paragraph("B-Factor Statistics", styles['Heading2']))
                story.append(Spacer(1, 6))
                
                b_data = [
                    ["Metric", "Value"],
                    ["Count", f"{len(b_vals):.0f}"],
                    ["Mean", f"{b_vals.mean():.3f}"],
                    ["Median", f"{b_vals.median():.3f}"],
                    ["P10", f"{np.percentile(b_vals, 10):.3f}"],
                    ["P90", f"{np.percentile(b_vals, 90):.3f}"]
                ]
                
                b_table = Table(b_data, colWidths=[2.5*inch, 2*inch])
                b_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(b_table)
                story.append(Spacer(1, 18))
        
       # Add plots
        hist_png, box_png, scatter_png, _ = bfactor_analytics_figures(oneline, fluid_name, eur_col)
        
        if hist_png:
            story.append(Paragraph("B-Factor Distribution", styles['Heading2']))
            story.append(Spacer(1, 6))
            story.append(Image(hist_png, width=5*inch, height=2.5*inch))
            story.append(Spacer(1, 12))
        
        # Create a new page for the three plots together
        story.append(PageBreak())
        story.append(Paragraph(f"{fluid_name} Analysis Charts", heading_style))
        story.append(Spacer(1, 12))
        
        # EUR vs B-Factor
        if scatter_png:
            story.append(Paragraph("EUR vs B-Factor", styles['Heading2']))
            story.append(Spacer(1, 6))
            story.append(Image(scatter_png, width=5.5*inch, height=2.6*inch))
            story.append(Spacer(1, 10))
        
        # Type curves — use the same curves already computed and scaled by _render_tw
        # so the PDF plot is guaranteed to match the UI plot exactly.
        _cached_curves = st.session_state.get(f"{fluid_name}_tc_curves", pd.DataFrame())
        _cached_lines  = st.session_state.get(f"{fluid_name}_tc_lines",  [])
        if _cached_curves.empty and monthly_sel is not None:
            # Fallback: recompute if session state was cleared (e.g. page reload before PDF)
            _cached_curves, _cached_lines = build_type_curves_and_lines(
                monthly_sel, fluid_name.lower(), oneline=oneline
            )
        if not _cached_curves.empty:
            curves = _cached_curves
            lines  = _cached_lines
            story.append(Paragraph("Type Curve", styles['Heading2']))
            story.append(Spacer(1, 6))
            fig = plot_type_curves(curves, lines, fluid_name.lower())
            tc_png = _save_fig(fig)
            story.append(Image(tc_png, width=5.5*inch, height=2.6*inch))
            story.append(Spacer(1, 8))

            # Arps parameters table for each type curve tier
            qi_col  = 'qi (per day)'
            di_col  = 'di (per month)'
            b_col   = 'b'
            dec_col = 'First-Year Decline (%)'
            if all(c in oneline.columns for c in [qi_col, b_col, di_col]):
                qi  = pd.to_numeric(oneline[qi_col],  errors='coerce').dropna()
                b_s = pd.to_numeric(oneline[b_col],   errors='coerce').dropna()
                di  = pd.to_numeric(oneline[di_col],  errors='coerce').dropna()
                dec = pd.to_numeric(oneline[dec_col], errors='coerce').dropna() if dec_col in oneline.columns else None

                qi_unit = "Mcf/day" if fluid_name == "Gas" else "bbl/day"

                tc_data = [["Parameter", "P90 Curve\n(Low)", "P50 Curve\n(Mid)", "P10 Curve\n(High)"]]
                tc_data.append([
                    f"qi ({qi_unit})",
                    f"{np.percentile(qi, 10):,.0f}",
                    f"{np.percentile(qi, 50):,.0f}",
                    f"{np.percentile(qi, 90):,.0f}",
                ])
                tc_data.append([
                    "b-factor",
                    f"{np.percentile(b_s, 10):.3f}",
                    f"{np.percentile(b_s, 50):.3f}",
                    f"{np.percentile(b_s, 90):.3f}",
                ])
                tc_data.append([
                    "Di (per month)",
                    f"{np.percentile(di, 90):.4f}",
                    f"{np.percentile(di, 50):.4f}",
                    f"{np.percentile(di, 10):.4f}",
                ])
                if dec is not None and not dec.empty:
                    tc_data.append([
                        "1st-Year Decline (%)",
                        f"{np.percentile(dec, 90):.1f}%",
                        f"{np.percentile(dec, 50):.1f}%",
                        f"{np.percentile(dec, 10):.1f}%",
                    ])
                tc_data.append(["Terminal Di (d_lim)", "0.00417/mo", "0.00417/mo", "0.00417/mo"])

                tc_table = Table(tc_data, colWidths=[2.0*inch, 1.4*inch, 1.4*inch, 1.4*inch])
                tc_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), _phase_color_rgb(fluid_name)),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('ALIGN', (0, 0), (0, -1), 'LEFT'),
                    ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                    ('TOPPADDING', (0, 0), (-1, -1), 5),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F5F5F5')]),
                ]))
                story.append(tc_table)
            story.append(Spacer(1, 10))
        
        # EUR summary table + probit plot
        if eur_col in oneline.columns:
            eurs = pd.to_numeric(oneline[eur_col], errors="coerce").dropna().astype(float).tolist()
            if eurs:
                unit = "Mbbl" if fluid_name != "Gas" else "MMcf"

                # EUR summary table
                eur_stats = compute_eur_stats(eurs)
                tbl_df = eur_summary_table(fluid_name, eur_stats, unit, int(norm_len))
                story.append(Paragraph("EUR Statistics", styles['Heading2']))
                story.append(Spacer(1, 4))
                col0, col1, col2 = tbl_df.columns
                eur_tbl_data = [[col0, col1, col2]] + [
                    [str(r[col0]),
                     f"{r[col1]:.2f}" if isinstance(r[col1], float) and np.isfinite(r[col1]) else (str(r[col1]) if not isinstance(r[col1], float) else ""),
                     f"{r[col2]:.2f}" if isinstance(r[col2], float) and np.isfinite(r[col2]) else (str(r[col2]) if not isinstance(r[col2], float) else "")]
                    for _, r in tbl_df.iterrows()
                ]
                eur_tbl = Table(eur_tbl_data, colWidths=[2.2*inch, 1.4*inch, 1.4*inch])
                eur_tbl.setStyle(TableStyle([
                    ('BACKGROUND',    (0, 0), (-1, 0), _phase_color_rgb(fluid_name)),
                    ('TEXTCOLOR',     (0, 0), (-1, 0), colors.whitesmoke),
                    ('FONTNAME',      (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('ALIGN',         (0, 0), (0, -1), 'LEFT'),
                    ('ALIGN',         (1, 0), (-1, -1), 'CENTER'),
                    ('FONTSIZE',      (0, 0), (-1, -1), 9),
                    ('TOPPADDING',    (0, 0), (-1, -1), 5),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                    ('GRID',          (0, 0), (-1, -1), 0.5, colors.black),
                    ('ROWBACKGROUNDS',(0, 1), (-1, -1), [colors.white, colors.HexColor('#F5F5F5')]),
                ]))
                story.append(eur_tbl)
                story.append(Spacer(1, 10))

                # Probit plot (EUR + EUR/ft side by side)
                story.append(Paragraph("Probit Plot", styles['Heading2']))
                story.append(Spacer(1, 6))
                fig = probit_plot(eurs, unit, f"{fluid_name} EUR Distribution",
                                  _phase_color(fluid_name), norm_len=int(norm_len))
                probit_png = _save_fig(fig)
                story.append(Image(probit_png, width=7.0*inch, height=3.0*inch))
                story.append(Spacer(1, 12))

        # ── Analog well oneline with Arps parameters ──
        story.append(PageBreak())
        story.append(Paragraph(f"{fluid_name} — Analog Wells & Arps Parameters", heading_style))
        story.append(Spacer(1, 8))
        story.append(Paragraph(
            f"The following {len(oneline)} wells were used to build the {fluid_name} type curve.",
            styles['Normal']))
        story.append(Spacer(1, 8))

        _well_cols_wanted = ['API10', 'WellName', 'LateralLength',
                              'qi (per day)', 'b', 'di (per month)', 'First-Year Decline (%)', eur_col]
        _well_cols = [c for c in _well_cols_wanted if c in oneline.columns]
        if _well_cols:
            _qi_unit = "Mcf/d" if fluid_name == "Gas" else "bbl/d"
            _hdr_map = {
                'API10': 'API/UWI', 'WellName': 'Well Name',
                'LateralLength': 'Lat Len (ft)',
                'qi (per day)': f'qi ({_qi_unit})', 'b': 'b',
                'di (per month)': 'Di (/mo)', 'First-Year Decline (%)': '1yr Dec (%)',
                eur_col: 'EUR'
            }
            _col_w = {
                'API10': 1.2*inch, 'WellName': 1.5*inch,
                'LateralLength': 0.7*inch, 'qi (per day)': 0.8*inch, 'b': 0.5*inch,
                'di (per month)': 0.6*inch, 'First-Year Decline (%)': 0.7*inch, eur_col: 0.7*inch
            }
            _headers = [_hdr_map.get(c, c) for c in _well_cols]
            _widths  = [_col_w.get(c, 0.9*inch) for c in _well_cols]

            _well_data = [_headers]
            for _, row in oneline[_well_cols].iterrows():
                _row_vals = []
                for c in _well_cols:
                    v = row[c]
                    if isinstance(v, float) and np.isfinite(v):
                        _row_vals.append(f"{v:,.2f}" if c not in ('b', 'di (per month)') else f"{v:.4f}")
                    else:
                        _row_vals.append(str(v) if pd.notna(v) else "")
                _well_data.append(_row_vals)

            _well_tbl = Table(_well_data, colWidths=_widths, repeatRows=1)
            _well_tbl.setStyle(TableStyle([
                ('BACKGROUND',    (0, 0), (-1, 0),  _phase_color_rgb(fluid_name)),
                ('TEXTCOLOR',     (0, 0), (-1, 0),  colors.whitesmoke),
                ('FONTNAME',      (0, 0), (-1, 0),  'Helvetica-Bold'),
                ('FONTSIZE',      (0, 0), (-1, -1), 7),
                ('ALIGN',         (0, 0), (-1, -1), 'CENTER'),
                ('ALIGN',         (0, 1), (1, -1),  'LEFT'),
                ('TOPPADDING',    (0, 0), (-1, -1), 3),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
                ('GRID',          (0, 0), (-1, -1), 0.4, colors.black),
                ('ROWBACKGROUNDS',(0, 1), (-1, -1), [colors.white, colors.HexColor('#F5F5F5')]),
            ]))
            story.append(_well_tbl)
        story.append(Spacer(1, 10))

        # ── Outlier detection: EUR vs Lateral Length ──
        _out = _make_outlier_fig(oneline, eur_col, fluid_name)
        if _out is not None:
            _of, _uf, _na, _nr, _nflag, _med, _unit = _out
            story.append(Spacer(1, 10))
            story.append(Paragraph(f"{fluid_name} — EUR Outlier Detection", styles['Heading2']))
            story.append(Spacer(1, 6))
            story.append(Image(_save_fig(_of), width=6.5*inch, height=3.2*inch))
            story.append(Spacer(1, 8))
            # Summary stats bar
            _stat_data = [[
                Paragraph("<b>Active wells</b><br/>" + str(_na),           styles['Normal']),
                Paragraph("<b>Removed wells</b><br/>" + str(_nr),          styles['Normal']),
                Paragraph(f"<b>Median EUR</b><br/>{_med:,.0f} {_unit}",   styles['Normal']),
                Paragraph(f"<b>Upper fence</b><br/>{_uf:,.0f} {_unit}",   styles['Normal']),
                Paragraph(f"<b>Flagged</b><br/>{_nflag} wells",            styles['Normal']),
            ]]
            _stat_tbl = Table(_stat_data, colWidths=[1.2*inch]*5)
            _stat_tbl.setStyle(TableStyle([
                ('BACKGROUND',    (0, 0), (-1, -1), colors.HexColor('#EFEFEF')),
                ('ALIGN',         (0, 0), (-1, -1), 'CENTER'),
                ('FONTSIZE',      (0, 0), (-1, -1), 8),
                ('TOPPADDING',    (0, 0), (-1, -1), 6),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('GRID',          (0, 0), (-1, -1), 0.5, colors.HexColor('#CCCCCC')),
            ]))
            story.append(_stat_tbl)
            story.append(Spacer(1, 10))

        story.append(PageBreak())
    
    # Add each fluid section
    add_fluid_section("Oil", "Oil_oneline", "Oil_monthly", "EUR (Mbbl)")
    add_fluid_section("Gas", "Gas_oneline", "Gas_monthly", "EUR (MMcf)")
    add_fluid_section("Water", "Water_oneline", "Water_monthly", "EUR (Mbbl water)")
    
    # Build PDF
    doc.build(story)
    return tmp.name

st.markdown("---")
st.header("Export Report")

_has_forecast = any(k in st.session_state for k in ["Oil_oneline", "Gas_oneline", "Water_oneline"])

if st.button("Generate PDF + PPTX Report", disabled=not _has_forecast):
    if not _has_forecast:
        st.error("Please run at least one forecast before generating a report.")
    else:
        st.session_state["_pdf_name_prompt"] = True

if not _has_forecast:
    st.info("Run at least one forecast (Oil, Gas, or Water) to enable report export.")

if st.session_state.get("_pdf_name_prompt"):
    tc_name = st.text_input(
        "Enter a name for this Type Curve report:",
        key="_pdf_tc_name",
        placeholder="e.g. Niobrara 2024 Q1",
    )
    if st.button("Confirm & Generate", key="_pdf_confirm"):
        try:
            name_clean = tc_name.strip()
            with st.spinner("Generating PDF report..."):
                pdf_path = generate_comprehensive_pdf(tc_name=name_clean)
            safe_name = "".join(c if c.isalnum() or c in " _-" else "_" for c in name_clean)
            pdf_file_name = f"{safe_name}.pdf" if safe_name else f"SE_Autoforecast_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            with open(pdf_path, "rb") as _f:
                st.session_state["_pdf_bytes"]     = _f.read()
                st.session_state["_pdf_file_name"] = pdf_file_name

            with st.spinner("Generating PPTX report..."):
                pptx_data = _build_pptx_data(name_clean)
                pptx_path = generate_tc_pptx(pptx_data, tempfile.gettempdir())
            if pptx_path and os.path.exists(pptx_path):
                with open(pptx_path, "rb") as _f:
                    st.session_state["_pptx_bytes"]     = _f.read()
                    st.session_state["_pptx_file_name"] = os.path.basename(pptx_path)
                st.success("PDF and PPTX reports ready for download.")
            else:
                st.session_state.pop("_pptx_bytes", None)
                st.success("PDF report ready for download. (PPTX export failed — see warning above)")

            st.session_state["_pdf_name_prompt"] = False
        except Exception as e:
            st.error(f"Error generating report: {str(e)}")
            st.exception(e)

# Persistent download buttons — survive rerenders via session state
_dl_col1, _dl_col2 = st.columns(2)
with _dl_col1:
    if st.session_state.get("_pdf_bytes"):
        st.download_button(
            label="Download PDF Report",
            data=st.session_state["_pdf_bytes"],
            file_name=st.session_state.get("_pdf_file_name", "SE_Report.pdf"),
            mime="application/pdf",
            key="_pdf_dl",
            use_container_width=True,
        )
with _dl_col2:
    if st.session_state.get("_pptx_bytes"):
        st.download_button(
            label="Download PPTX Report",
            data=st.session_state["_pptx_bytes"],
            file_name=st.session_state.get("_pptx_file_name", "SE_Report.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="_pptx_dl",
            use_container_width=True,
        )
